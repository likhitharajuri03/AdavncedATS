import os
import io
import re
from typing import Tuple, Optional
import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Advanced document processor supporting multiple formats
    with resume detection capability
    """
    
    RESUME_KEYWORDS = [
        'experience', 'education', 'skills', 'work history',
        'employment', 'qualification', 'objective', 'summary',
        'projects', 'certifications', 'achievements', 'responsibilities',
        'profile', 'career', 'contact', 'email', 'phone', 'address',
        'bachelor', 'master', 'degree', 'university', 'college',
        'internship', 'training', 'volunteer', 'languages', 'hobbies'
    ]
    
    RESUME_PATTERNS = [
        r'\b(?:java|python|javascript|react|node|sql|html|css)\b',
        r'\b(?:manager|developer|engineer|analyst|designer|consultant)\b',
        r'\b(?:2019|2020|2021|2022|2023|2024|2025)\b',
        r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
        r'\b(?:\+\d{1,3}[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b'
    ]
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.doc', '.ppt', '.pptx', '.jpg', '.jpeg', '.png']
    
    def process_document(self, file_path: str, file_bytes: bytes = None) -> Tuple[str, bool, float]:
        """
        Process document and return extracted text, is_resume flag, and confidence score
        
        Args:
            file_path: Path to the file or filename
            file_bytes: Bytes content of the file
            
        Returns:
            Tuple of (extracted_text, is_resume, confidence_score)
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {ext}")
        
        try:
            if ext == '.pdf':
                text = self._extract_from_pdf(file_bytes)
            elif ext in ['.docx', '.doc']:
                text = self._extract_from_docx(file_bytes)
            elif ext in ['.ppt', '.pptx']:
                text = self._extract_from_pptx(file_bytes)
            elif ext in ['.jpg', '.jpeg', '.png']:
                text = self._extract_from_image(file_bytes)
            else:
                raise ValueError(f"Unexpected format: {ext}")
            
            # Detect if document is a resume
            is_resume, confidence = self._detect_resume(text)
            
            logger.info(f"Processed {ext} document. Is Resume: {is_resume}, Confidence: {confidence:.2f}")
            
            return text, is_resume, confidence
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            raise
    
    def _extract_from_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF using multiple methods for better accuracy"""
        text = ""
        
        # Method 1: PyPDF2
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        except Exception as e:
            logger.warning(f"PyPDF2 extraction failed: {str(e)}")
        
        # Method 2: pdfplumber (better for complex layouts)
        if len(text.strip()) < 100:
            try:
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                logger.warning(f"pdfplumber extraction failed: {str(e)}")
        
        # Method 3: OCR for scanned PDFs
        if len(text.strip()) < 100:
            try:
                images = convert_from_bytes(file_bytes)
                for image in images:
                    text += pytesseract.image_to_string(image) + "\n"
            except Exception as e:
                logger.warning(f"OCR extraction failed: {str(e)}")
        
        return text.strip()
    
    def _extract_from_docx(self, file_bytes: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(io.BytesIO(file_bytes))
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += " " + cell.text
            
            return text.strip()
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            raise
    
    def _extract_from_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            raise
    
    def _extract_from_image(self, file_bytes: bytes) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(io.BytesIO(file_bytes))
            
            # Preprocess image for better OCR
            image = image.convert('L')  # Convert to grayscale
            
            # Extract text using Tesseract
            text = pytesseract.image_to_string(image)
            
            return text.strip()
        except Exception as e:
            logger.error(f"Image extraction failed: {str(e)}")
            raise
    
    def _detect_resume(self, text: str) -> Tuple[bool, float]:
        """
        Detect if the extracted text is from a resume
        Returns (is_resume, confidence_score)
        """
        if not text or len(text.strip()) < 50:
            return False, 0.0
        
        text_lower = text.lower()
        word_count = len(text.split())
        
        # Score based on keyword presence
        keyword_score = 0
        keywords_found = []
        for keyword in self.RESUME_KEYWORDS:
            if keyword in text_lower:
                keyword_score += 1
                keywords_found.append(keyword)
        
        # Score based on pattern matching
        pattern_score = 0
        for pattern in self.RESUME_PATTERNS:
            if re.search(pattern, text, re.IGNORECASE):
                pattern_score += 1
        
        # Calculate confidence
        keyword_confidence = min(keyword_score / 10, 1.0)  # Normalize to 0-1
        pattern_confidence = min(pattern_score / 3, 1.0)
        
        # Word count factor (resumes typically 200-2000 words)
        word_factor = 1.0 if 200 <= word_count <= 2000 else 0.5
        
        # Combined confidence score
        confidence = (keyword_confidence * 0.6 + pattern_confidence * 0.3 + word_factor * 0.1)
        
        # Threshold for resume detection
        is_resume = confidence >= 0.4
        
        logger.info(f"Resume Detection - Keywords: {len(keywords_found)}, Patterns: {pattern_score}, Words: {word_count}")
        
        return is_resume, confidence
    
    def get_document_metadata(self, text: str) -> dict:
        """Extract basic metadata from resume text"""
        metadata = {
            'word_count': len(text.split()),
            'char_count': len(text),
            'email': None,
            'phone': None,
            'has_skills_section': False,
            'has_experience_section': False,
            'has_education_section': False
        }
        
        # Extract email
        email_match = re.search(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text)
        if email_match:
            metadata['email'] = email_match.group()
        
        # Extract phone
        phone_match = re.search(r'\b(?:\+\d{1,3}[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b', text)
        if phone_match:
            metadata['phone'] = phone_match.group()
        
        # Check for sections
        text_lower = text.lower()
        metadata['has_skills_section'] = any(keyword in text_lower for keyword in ['skills', 'technical skills', 'core competencies'])
        metadata['has_experience_section'] = any(keyword in text_lower for keyword in ['experience', 'work history', 'employment'])
        metadata['has_education_section'] = any(keyword in text_lower for keyword in ['education', 'academic', 'qualification'])
        
        return metadata